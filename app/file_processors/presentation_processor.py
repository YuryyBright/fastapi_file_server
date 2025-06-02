from pptx import Presentation
from app.file_processors.base_processor import FileProcessor


class PresentationProcessor(FileProcessor):
    """
    Processor for PowerPoint presentation files.
    Supports both `.pptx` files for text extraction and search functionality.
    """

    def __init__(self, file_path: str):
        """
        Initialize the PowerPoint processor with the file path.
        :param file_path: Path to the PowerPoint presentation file.
        """
        self.file_path = file_path

    def search(self, keyword: str, exact_match: bool) -> list[str]:
        """
        Search for a keyword in a PowerPoint presentation.
        :param keyword: The keyword to search for within the presentation.
        :return: A list of paragraphs containing the keyword.
        """
        results = []
        presentation = Presentation(self.file_path)

        # Iterate through each slide in the presentation
        for slide in presentation.slides:
            for shape in slide.shapes:
                # If the shape has a text frame, check its paragraphs
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if keyword in paragraph.text:
                            results.append(paragraph.text)  # Add paragraph text if keyword is found
        return results

    def read(self) -> str:
        """
        Read the entire content of the PowerPoint presentation and return it as a single string.
        :return: The text content of the entire PowerPoint presentation.
        """
        full_text = ""
        presentation = Presentation(self.file_path)

        # Iterate through each slide and each shape in the slide
        for slide in presentation.slides:
            for shape in slide.shapes:
                # If the shape has a text frame, extract its paragraphs
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        full_text += paragraph.text + "\n"  # Append paragraph text to the full_text
        return full_text.strip()  # Return the content as a single string
